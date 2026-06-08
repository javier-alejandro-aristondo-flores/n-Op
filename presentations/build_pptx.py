'''Convert files/slides.md into a dark-themed .pptx via python-pptx.

Source of truth is slides.md. This script does no rewording. Each level-2
heading whose text starts with "Slide" becomes one slide. Inside a slide:

- A leading blockquote (lines starting with `>`) becomes an italic lede
  beneath the title.
- Fenced code blocks become monospace panels with preserved whitespace.
- Other prose lines become a body text block (split on blank lines into
  paragraphs).
- Horizontal rules `---` are slide separators (already implied by the
  heading split; treated as no-ops).
'''


import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Emu, Inches, Pt


SESSION_FILES_DIR: Path = Path(__file__).resolve().parent
SLIDES_MD_PATH: Path = SESSION_FILES_DIR / "slides.md"
OUTPUT_PPTX_PATH: Path = SESSION_FILES_DIR / "n-op-physics-catchup.pptx"


SLIDE_WIDTH_IN: float = 13.333
SLIDE_HEIGHT_IN: float = 7.5

BACKGROUND_HEX: str = "1A1A1A"
TITLE_HEX: str = "F0F0F0"
BODY_HEX: str = "D0D0D0"
LEDE_HEX: str = "B0B0B0"
CODE_BG_HEX: str = "252525"
CODE_FG_HEX: str = "9CDCFE"
ACCENT_HEX: str = "B5CEA8"
FOOTER_HEX: str = "6A6A6A"

TITLE_FONT_NAME: str = "Calibri"
BODY_FONT_NAME: str = "Calibri"
CODE_FONT_NAME: str = "Consolas"

TITLE_FONT_SIZE_PT: int = 32
LEDE_FONT_SIZE_PT: int = 18
BODY_FONT_SIZE_PT: int = 16
CODE_FONT_SIZE_PT: int = 13
FOOTER_FONT_SIZE_PT: int = 10

ACCENT_TOKENS: tuple[str, ...] = (
    "RESIDUAL",
    "ASSERT",
    "==",
    "PULLBACK_CHECK",
    "VALID_STATE_PREDICATE",
)


def _hex_to_rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str)


def _set_solid_fill(shape: object, hex_str: str) -> None:
    fill = shape.fill  # type: ignore[attr-defined]
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(hex_str)


def _no_line(shape: object) -> None:
    line = shape.line  # type: ignore[attr-defined]
    line.fill.background()


class SlideBlock:
    '''A typed body block: prose paragraph or code panel.'''

    __slots__ = ("kind", "text")

    def __init__(self, kind: str, text: str) -> None:
        self.kind: str = kind
        self.text: str = text


class ParsedSlide:
    '''One parsed slide: title, optional lede, ordered body blocks.'''

    __slots__ = ("number", "title", "lede", "blocks")

    def __init__(
        self,
        number: int,
        title: str,
        lede: str,
        blocks: list[SlideBlock],
    ) -> None:
        self.number: int = number
        self.title: str = title
        self.lede: str = lede
        self.blocks: list[SlideBlock] = blocks


SLIDE_HEADING_PATTERN: re.Pattern[str] = re.compile(
    r"^##\s+Slide\s+(\d+)\s*[—-]\s*(.+?)\s*$"
)


def Parse_Slides_Markdown(markdown_text: str) -> list[ParsedSlide]:
    '''Split markdown into ParsedSlide list.'''
    lines: list[str] = markdown_text.splitlines()
    parsed_slides: list[ParsedSlide] = []
    current_number: int = 0
    current_title: str = ""
    current_lede_lines: list[str] = []
    current_blocks: list[SlideBlock] = []
    current_prose_buffer: list[str] = []
    in_code: bool = False
    current_code_buffer: list[str] = []
    seen_first_heading: bool = False

    def Flush_Prose() -> None:
        nonlocal current_prose_buffer
        joined: str = "\n".join(current_prose_buffer).strip()
        if joined:
            current_blocks.append(SlideBlock("prose", joined))
        current_prose_buffer = []

    def Flush_Slide() -> None:
        nonlocal current_lede_lines, current_blocks
        if not seen_first_heading:
            return
        Flush_Prose()
        lede_text: str = " ".join(
            line.lstrip("> ").strip() for line in current_lede_lines
        ).strip()
        lede_text = lede_text.strip("*").strip()
        parsed_slides.append(
            ParsedSlide(
                number=current_number,
                title=current_title,
                lede=lede_text,
                blocks=list(current_blocks),
            )
        )
        current_lede_lines = []
        current_blocks = []

    for raw_line in lines:
        if in_code:
            if raw_line.strip().startswith("```"):
                code_text: str = "\n".join(current_code_buffer)
                current_blocks.append(SlideBlock("code", code_text))
                current_code_buffer = []
                in_code = False
            else:
                current_code_buffer.append(raw_line)
            continue
        heading_match: re.Match[str] | None = SLIDE_HEADING_PATTERN.match(raw_line)
        if heading_match is not None:
            Flush_Slide()
            seen_first_heading = True
            current_number = int(heading_match.group(1))
            current_title = heading_match.group(2).strip()
            current_lede_lines = []
            current_blocks = []
            current_prose_buffer = []
            continue
        if not seen_first_heading:
            continue
        if raw_line.strip().startswith("```"):
            Flush_Prose()
            in_code = True
            current_code_buffer = []
            continue
        if raw_line.strip() == "---":
            continue
        is_blockquote: bool = raw_line.lstrip().startswith(">")
        prose_has_content: bool = any(line.strip() for line in current_prose_buffer)
        if is_blockquote and not current_blocks and not prose_has_content:
            current_lede_lines.append(raw_line)
            continue
        current_prose_buffer.append(raw_line)

    Flush_Slide()
    return parsed_slides


def Apply_Background(slide: object, hex_color: str) -> None:
    background = slide.background  # type: ignore[attr-defined]
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(hex_color)


def Add_Title(
    slide: object,
    slide_number: int,
    title_text: str,
) -> None:
    box = slide.shapes.add_textbox(  # type: ignore[attr-defined]
        Inches(0.5), Inches(0.3), Inches(SLIDE_WIDTH_IN - 1.0), Inches(0.7)
    )
    text_frame = box.text_frame
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    number_run = paragraph.add_run()
    number_run.text = f"{slide_number:02d}   "
    number_run.font.name = TITLE_FONT_NAME
    number_run.font.size = Pt(TITLE_FONT_SIZE_PT)
    number_run.font.bold = True
    number_run.font.color.rgb = _hex_to_rgb(FOOTER_HEX)
    title_run = paragraph.add_run()
    title_run.text = title_text
    title_run.font.name = TITLE_FONT_NAME
    title_run.font.size = Pt(TITLE_FONT_SIZE_PT)
    title_run.font.bold = True
    title_run.font.color.rgb = _hex_to_rgb(TITLE_HEX)


def _Estimate_Prose_Height(text: str, width_in: float, font_size_pt: float) -> float:
    chars_per_in: float = 90.0 / 6.4 * (16.0 / font_size_pt)
    line_height_in: float = font_size_pt * 0.020
    chars_per_line: float = max(20.0, chars_per_in * width_in)
    paragraphs: list[str] = [c.strip() for c in re.split(r"\n\s*\n", text) if c.strip()]
    total_in: float = 0.0
    for chunk in paragraphs:
        cleaned = chunk.replace("\n", " ")
        line_count = max(1, int(len(cleaned) / chars_per_line) + 1)
        total_in += line_height_in * line_count + 0.08
    return max(0.3, total_in)


def _Estimate_Code_Height(text: str, font_size_pt: float) -> float:
    line_count: int = len(text.split("\n"))
    line_height_in: float = font_size_pt * 0.0185
    return line_height_in * line_count + 0.3


def _Compute_Scale(parsed_slide: "ParsedSlide", body_width_in: float) -> float:
    available_in: float = SLIDE_HEIGHT_IN - 1.15 - 0.5
    if parsed_slide.lede:
        available_in -= 0.75
    total_in: float = 0.0
    fixed_in: float = 0.0
    for block in parsed_slide.blocks:
        if block.kind == "code":
            total_in += _Estimate_Code_Height(block.text, CODE_FONT_SIZE_PT)
            fixed_in += 0.12
        else:
            total_in += _Estimate_Prose_Height(block.text, body_width_in, BODY_FONT_SIZE_PT)
            fixed_in += 0.10
    if total_in + fixed_in <= available_in:
        return 1.0
    variable: float = max(0.1, total_in)
    headroom: float = max(0.5, available_in - fixed_in)
    return max(0.4, headroom / variable)


def Add_Lede(slide: object, lede_text: str, top_in: float) -> float:
    if not lede_text:
        return top_in
    box = slide.shapes.add_textbox(  # type: ignore[attr-defined]
        Inches(0.5),
        Inches(top_in),
        Inches(SLIDE_WIDTH_IN - 1.0),
        Inches(0.7),
    )
    text_frame = box.text_frame
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = lede_text
    run.font.name = BODY_FONT_NAME
    run.font.size = Pt(LEDE_FONT_SIZE_PT)
    run.font.italic = True
    run.font.color.rgb = _hex_to_rgb(LEDE_HEX)
    return top_in + 0.7


def _Add_Prose_Run(paragraph: object, text_segment: str, is_bold: bool, font_size_pt: float) -> None:
    run = paragraph.add_run()  # type: ignore[attr-defined]
    run.text = text_segment
    run.font.name = BODY_FONT_NAME
    run.font.size = Pt(font_size_pt)
    run.font.color.rgb = _hex_to_rgb(BODY_HEX)
    if is_bold:
        run.font.bold = True


def _Render_Inline_Markdown(paragraph: object, text: str, font_size_pt: float) -> None:
    pattern: re.Pattern[str] = re.compile(r"(\*\*[^*]+\*\*|`[^`]+`|\*[^*]+\*)")
    cursor: int = 0
    for match in pattern.finditer(text):
        if match.start() > cursor:
            _Add_Prose_Run(paragraph, text[cursor : match.start()], False, font_size_pt)
        token: str = match.group(0)
        if token.startswith("**"):
            _Add_Prose_Run(paragraph, token[2:-2], True, font_size_pt)
        elif token.startswith("`"):
            run = paragraph.add_run()  # type: ignore[attr-defined]
            run.text = token[1:-1]
            run.font.name = CODE_FONT_NAME
            run.font.size = Pt(font_size_pt - 1)
            run.font.color.rgb = _hex_to_rgb(CODE_FG_HEX)
        else:
            run = paragraph.add_run()  # type: ignore[attr-defined]
            run.text = token[1:-1]
            run.font.name = BODY_FONT_NAME
            run.font.size = Pt(font_size_pt)
            run.font.italic = True
            run.font.color.rgb = _hex_to_rgb(BODY_HEX)
        cursor = match.end()
    if cursor < len(text):
        _Add_Prose_Run(paragraph, text[cursor:], False, font_size_pt)


def Add_Prose_Block(
    slide: object,
    text: str,
    top_in: float,
    width_in: float,
    scale: float,
) -> float:
    body_size: float = BODY_FONT_SIZE_PT * scale
    paragraphs_text: list[str] = [
        chunk.strip() for chunk in re.split(r"\n\s*\n", text) if chunk.strip()
    ]
    estimated_height_in: float = _Estimate_Prose_Height(text, width_in, body_size)
    box = slide.shapes.add_textbox(  # type: ignore[attr-defined]
        Inches(0.5),
        Inches(top_in),
        Inches(width_in),
        Inches(estimated_height_in),
    )
    text_frame = box.text_frame
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    text_frame.word_wrap = True
    for index, chunk in enumerate(paragraphs_text):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.space_after = Pt(4 * scale)
        cleaned: str = chunk.replace("\n", " ").strip()
        is_bullet: bool = cleaned.startswith("- ") or cleaned.startswith("* ")
        if is_bullet:
            cleaned = cleaned[2:]
            bullet_run = paragraph.add_run()
            bullet_run.text = "• "
            bullet_run.font.name = BODY_FONT_NAME
            bullet_run.font.size = Pt(body_size)
            bullet_run.font.color.rgb = _hex_to_rgb(ACCENT_HEX)
        _Render_Inline_Markdown(paragraph, cleaned, body_size)
    return top_in + estimated_height_in + 0.1


def _Token_Run(paragraph: object, text_segment: str, color_hex: str, font_size_pt: float) -> None:
    run = paragraph.add_run()  # type: ignore[attr-defined]
    run.text = text_segment
    run.font.name = CODE_FONT_NAME
    run.font.size = Pt(font_size_pt)
    run.font.color.rgb = _hex_to_rgb(color_hex)


def _Render_Code_Line(paragraph: object, line_text: str, font_size_pt: float) -> None:
    if not line_text:
        _Token_Run(paragraph, "", CODE_FG_HEX, font_size_pt)
        return
    accent_pattern: re.Pattern[str] = re.compile(
        r"\b(" + "|".join(re.escape(t) for t in ACCENT_TOKENS) + r")\b|=="
    )
    cursor: int = 0
    for match in accent_pattern.finditer(line_text):
        if match.start() > cursor:
            _Token_Run(paragraph, line_text[cursor : match.start()], CODE_FG_HEX, font_size_pt)
        _Token_Run(paragraph, match.group(0), ACCENT_HEX, font_size_pt)
        cursor = match.end()
    if cursor < len(line_text):
        _Token_Run(paragraph, line_text[cursor:], CODE_FG_HEX, font_size_pt)


def Add_Code_Block(
    slide: object,
    code_text: str,
    top_in: float,
    width_in: float,
    scale: float,
) -> float:
    code_size: float = CODE_FONT_SIZE_PT * scale
    code_lines: list[str] = code_text.split("\n")
    estimated_height_in: float = _Estimate_Code_Height(code_text, code_size)
    panel = slide.shapes.add_shape(  # type: ignore[attr-defined]
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5),
        Inches(top_in),
        Inches(width_in),
        Inches(estimated_height_in),
    )
    _set_solid_fill(panel, CODE_BG_HEX)
    _no_line(panel)
    panel.text_frame.margin_left = Inches(0.18)
    panel.text_frame.margin_right = Inches(0.18)
    panel.text_frame.margin_top = Inches(0.10)
    panel.text_frame.margin_bottom = Inches(0.10)
    panel.text_frame.word_wrap = False
    for index, code_line in enumerate(code_lines):
        paragraph = (
            panel.text_frame.paragraphs[0]
            if index == 0
            else panel.text_frame.add_paragraph()
        )
        paragraph.space_after = Pt(0)
        paragraph.space_before = Pt(0)
        paragraph.line_spacing = 1.05
        _Render_Code_Line(paragraph, code_line.rstrip("\n"), code_size)
    return top_in + estimated_height_in + 0.12


def Add_Footer(slide: object, slide_number: int, total_slides: int) -> None:
    box = slide.shapes.add_textbox(  # type: ignore[attr-defined]
        Inches(SLIDE_WIDTH_IN - 1.6),
        Inches(SLIDE_HEIGHT_IN - 0.45),
        Inches(1.3),
        Inches(0.3),
    )
    text_frame = box.text_frame
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = 2
    run = paragraph.add_run()
    run.text = f"{slide_number} / {total_slides}"
    run.font.name = BODY_FONT_NAME
    run.font.size = Pt(FOOTER_FONT_SIZE_PT)
    run.font.color.rgb = _hex_to_rgb(FOOTER_HEX)


def Render_Slide(
    presentation: object,
    parsed_slide: ParsedSlide,
    total_slides: int,
) -> None:
    blank_layout = presentation.slide_layouts[6]  # type: ignore[attr-defined]
    slide = presentation.slides.add_slide(blank_layout)  # type: ignore[attr-defined]
    Apply_Background(slide, BACKGROUND_HEX)
    Add_Title(slide, parsed_slide.number, parsed_slide.title)
    cursor_top_in: float = 1.15
    if parsed_slide.lede:
        cursor_top_in = Add_Lede(slide, parsed_slide.lede, cursor_top_in) + 0.05
    body_width_in: float = SLIDE_WIDTH_IN - 1.0
    scale: float = _Compute_Scale(parsed_slide, body_width_in)
    for block in parsed_slide.blocks:
        if block.kind == "code":
            cursor_top_in = Add_Code_Block(
                slide, block.text, cursor_top_in, body_width_in, scale
            )
        else:
            cursor_top_in = Add_Prose_Block(
                slide, block.text, cursor_top_in, body_width_in, scale
            )
    Add_Footer(slide, parsed_slide.number, total_slides)


def Build_Presentation(parsed_slides: list[ParsedSlide]) -> Presentation:
    presentation: Presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_WIDTH_IN)
    presentation.slide_height = Inches(SLIDE_HEIGHT_IN)
    total: int = len(parsed_slides)
    for parsed_slide in parsed_slides:
        Render_Slide(presentation, parsed_slide, total)
    return presentation


def Main() -> None:
    markdown_text: str = SLIDES_MD_PATH.read_text(encoding="utf-8")
    parsed_slides: list[ParsedSlide] = Parse_Slides_Markdown(markdown_text)
    presentation: Presentation = Build_Presentation(parsed_slides)
    presentation.save(str(OUTPUT_PPTX_PATH))
    print(f"wrote {OUTPUT_PPTX_PATH} with {len(parsed_slides)} slides")


if __name__ == "__main__":
    Main()
