#!/usr/bin/env python3
"""Build the dark-mode deck from the talk script and the light figures.

  python3 build_pptx.py

Reads   2026-07-22-talk-script.md  and  figures/fig-*.svg
Writes  figures/dark/fig-*.svg     (palette-inverted)
        build/fig-*.png            (rasterised, 3x)
        n-op-talk.pptx             (16:9, dark, speaker notes attached)

The light SVGs stay the source of truth. Dark variants are generated, never hand-edited.
"""
import re, pathlib, io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
import cairosvg

HERE = pathlib.Path(__file__).parent
FIGS = HERE / "figures"
DARK = FIGS / "dark"
BUILD = HERE / "build"
SCRIPT = HERE / "2026-07-22-talk-script.md"
OUT = HERE / "n-op-talk.pptx"

BG = "#1c1c1c"
INK = RGBColor(0xE6, 0xE6, 0xE6)
DIM = RGBColor(0x8A, 0x8A, 0x8A)

# light -> dark. No target collides with a source, so sequential replace is safe.
PALETTE = [
    ("#000000", "#e6e6e6"),   # ink
    ("#ffffff", "#1c1c1c"),   # paper fills
    ("#eeeeee", "#303030"),   # light fill
    ("#cccccc", "#4a4a4a"),   # hairlines / mid fill
    ("#999999", "#8a8a8a"),   # grey marks, hatch
]
SCALE = 3.0                    # raster factor


def darken(svg: str) -> str:
    for light, dark in PALETTE:
        svg = svg.replace(light, dark)
    # <text> carries no explicit fill and would otherwise inherit SVG's default black.
    # No shape in the set relies on default fill, so setting it at the root hits text only.
    svg = svg.replace("<svg xmlns=", '<svg fill="#e6e6e6" xmlns=', 1)
    # self-contained background, matching the slide
    return svg.replace("</defs>", f'</defs>\n  <rect x="-9999" y="-9999" width="99999" '
                                  f'height="99999" fill="{BG}"/>', 1)


def parse_script():
    """-> [(number, title, figure-stem or None, [note lines])]"""
    slides, cur = [], None
    for line in SCRIPT.read_text(encoding="utf-8").splitlines():
        h = re.match(r"^## (\d+) · (.+?)(?: · `(fig-\d+)`)?\s*$", line)
        if h:
            if cur: slides.append(cur)
            cur = (int(h.group(1)), h.group(2).strip(), h.group(3), [])
            continue
        if line.startswith("## "):          # a non-slide section ends accumulation
            if cur: slides.append(cur)
            cur = None
            continue
        if cur is None or line.startswith("---"):
            continue
        b = re.match(r"^\s*[-*]\s+(.*)$", line)
        if b:
            cur[3].append(b.group(1).rstrip())
        elif line.strip() and cur[3] and line.startswith("  "):
            cur[3][-1] += " " + line.strip()          # continuation of a wrapped bullet
    if cur: slides.append(cur)
    return slides          # non-slide sections already ended accumulation above


def strip_md(s: str) -> str:
    s = re.sub(r"\*\*(.+?)\*\*", r"\1", s)
    s = re.sub(r"\*(.+?)\*", r"\1", s)
    return s.replace("`", "")


def main():
    DARK.mkdir(exist_ok=True); BUILD.mkdir(exist_ok=True)

    stems, light = {}, set()
    for src in sorted(FIGS.glob("fig-*.svg")):
        dsvg = DARK / src.name
        dsvg.write_text(darken(src.read_text(encoding="utf-8")), encoding="utf-8")
        png = BUILD / (src.stem + ".png")
        cairosvg.svg2png(url=str(dsvg), write_to=str(png), scale=SCALE)
        stems[src.name.split("-")[0] + "-" + src.name.split("-")[1]] = png
    # Raster figures pass through untouched. They carry their own light ground, so the
    # slide under them goes light too rather than framing white in black.
    for src in sorted(FIGS.glob("fig-*.png")):
        key = src.name.split("-")[0] + "-" + src.name.split("-")[1]
        stems[key] = src
        light.add(key)
    print(f"  {len(stems)} figures ({len(light)} raster, passed through)")

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    deck = parse_script()
    ends = (deck[0][0], deck[-1][0])          # title and close, whatever they are numbered
    for num, title, fig, notes in deck:
        s = prs.slides.add_slide(blank)
        on_light = fig in light
        bg = s.background.fill; bg.solid()
        bg.fore_color.rgb = RGBColor(0xF4, 0xF4, 0xF4) if on_light else RGBColor(0x1C, 0x1C, 0x1C)

        tb = s.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.6))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run(); r.text = strip_md(title)
        r.font.size = Pt(20 if num not in ends else 30)
        r.font.color.rgb = RGBColor(0x1C, 0x1C, 0x1C) if on_light else INK
        r.font.name = "Georgia"

        png = stems.get(fig) if fig else None
        if png and png.exists():
            from PIL import Image
            with Image.open(png) as im: iw, ih = im.size
            avail_w, avail_h = Inches(11.4), Inches(5.7)
            k = min(avail_w / iw, avail_h / ih)
            w, h = int(iw * k), int(ih * k)
            s.shapes.add_picture(str(png), int((prs.slide_width - w) / 2),
                                 Inches(1.35) + int((avail_h - h) / 2), w, h)
        elif num in ends:
            sub = s.shapes.add_textbox(Inches(0.6), Inches(3.2), Inches(12.1), Inches(1.2))
            q = sub.text_frame.paragraphs[0]
            rr = q.add_run()
            rr.text = ("Javier Flores  ·  n-Op" if num == ends[0]
                       else "Verifying is cheaper than solving.")
            rr.font.size = Pt(18); rr.font.color.rgb = DIM; rr.font.name = "Georgia"

        if notes:
            s.notes_slide.notes_text_frame.text = "\n".join("• " + strip_md(n) for n in notes)

    prs.save(OUT)
    print(f"  wrote {OUT.name} — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")


if __name__ == "__main__":
    main()
